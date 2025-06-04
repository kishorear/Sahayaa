"""
Instruction Document Processor
Scans for instruction files (.txt, .pdf, .docx, .pptx, .xlsx), converts to Markdown,
generates embeddings, and stores in Qdrant for similarity search.
"""

import os
import sys
import logging
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
import json
import hashlib
from datetime import datetime

# Document processing imports
try:
    from docx import Document
    from openpyxl import load_workbook
    from PyPDF2 import PdfReader
    from pptx import Presentation
    OFFICE_SUPPORT = True
except ImportError as e:
    logging.warning(f"Office document support limited: {e}")
    OFFICE_SUPPORT = False

# Vector and embedding imports
from qdrant_client import QdrantClient
from qdrant_client.http import models
from qdrant_client.http.models import Distance, VectorParams, PointStruct
import openai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class DocumentConverter:
    """Converts various document formats to Markdown text."""
    
    def __init__(self):
        self.supported_extensions = {'.txt', '.pdf', '.docx', '.pptx', '.xlsx'}
    
    def convert_to_markdown(self, file_path: Path) -> str:
        """Convert document to Markdown format."""
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.txt':
                return self._convert_txt(file_path)
            elif extension == '.pdf':
                return self._convert_pdf(file_path)
            elif extension == '.docx':
                return self._convert_docx(file_path)
            elif extension == '.pptx':
                return self._convert_pptx(file_path)
            elif extension == '.xlsx':
                return self._convert_xlsx(file_path)
            else:
                raise ValueError(f"Unsupported file format: {extension}")
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            raise
    
    def _convert_txt(self, file_path: Path) -> str:
        """Convert plain text file to Markdown."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read().strip()
        
        # Add filename as header
        filename = file_path.stem
        return f"# {filename}\n\n{content}"
    
    def _convert_pdf(self, file_path: Path) -> str:
        """Convert PDF to Markdown."""
        if not OFFICE_SUPPORT:
            raise ImportError("PDF processing requires PyPDF2")
        
        try:
            reader = PdfReader(str(file_path))
            text_content = []
            
            # Extract text from each page
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    text_content.append(f"## Page {i + 1}\n\n{page_text}")
            
            filename = file_path.stem
            content = "\n\n".join(text_content)
            return f"# {filename}\n\n{content}"
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            raise
    
    def _convert_docx(self, file_path: Path) -> str:
        """Convert Word document to Markdown."""
        if not OFFICE_SUPPORT:
            raise ImportError("DOCX processing requires python-docx")
        
        try:
            doc = Document(str(file_path))
            content_parts = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # Simple formatting detection
                    if paragraph.style.name.startswith('Heading'):
                        level = min(int(paragraph.style.name.split()[-1]) if paragraph.style.name.split()[-1].isdigit() else 1, 6)
                        content_parts.append(f"{'#' * level} {text}")
                    else:
                        content_parts.append(text)
            
            # Process tables
            for table in doc.tables:
                table_md = self._table_to_markdown(table)
                if table_md:
                    content_parts.append(table_md)
            
            filename = file_path.stem
            content = "\n\n".join(content_parts)
            return f"# {filename}\n\n{content}"
        except Exception as e:
            logger.error(f"Error reading DOCX {file_path}: {e}")
            raise
    
    def _convert_pptx(self, file_path: Path) -> str:
        """Convert PowerPoint presentation to Markdown."""
        if not OFFICE_SUPPORT:
            raise ImportError("PPTX processing requires python-pptx")
        
        try:
            prs = Presentation(str(file_path))
            slides_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slide_content = "\n".join(slide_text)
                    slides_content.append(f"## Slide {i + 1}\n\n{slide_content}")
            
            filename = file_path.stem
            content = "\n\n".join(slides_content)
            return f"# {filename}\n\n{content}"
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path}: {e}")
            raise
    
    def _convert_xlsx(self, file_path: Path) -> str:
        """Convert Excel spreadsheet to Markdown."""
        if not OFFICE_SUPPORT:
            raise ImportError("XLSX processing requires openpyxl")
        
        try:
            workbook = load_workbook(str(file_path), data_only=True)
            sheets_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                
                # Read all rows with data
                for row in sheet.iter_rows(values_only=True):
                    if any(cell for cell in row if cell is not None):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        sheet_data.append("| " + " | ".join(row_data) + " |")
                
                if sheet_data:
                    # Add table header separator
                    if len(sheet_data) > 1:
                        header_separator = "| " + " | ".join(["---"] * len(sheet_data[0].split("|")[1:-1])) + " |"
                        sheet_data.insert(1, header_separator)
                    
                    sheet_content = f"## {sheet_name}\n\n" + "\n".join(sheet_data)
                    sheets_content.append(sheet_content)
            
            filename = file_path.stem
            content = "\n\n".join(sheets_content)
            return f"# {filename}\n\n{content}"
        except Exception as e:
            logger.error(f"Error reading XLSX {file_path}: {e}")
            raise
    
    def _table_to_markdown(self, table) -> str:
        """Convert Word table to Markdown format."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):  # Only add non-empty rows
                rows.append("| " + " | ".join(cells) + " |")
        
        if len(rows) > 1:
            # Add header separator
            header_separator = "| " + " | ".join(["---"] * len(rows[0].split("|")[1:-1])) + " |"
            rows.insert(1, header_separator)
        
        return "\n".join(rows) if rows else ""

class EmbeddingGenerator:
    """Generates embeddings using OpenAI or local models."""
    
    def __init__(self):
        self.client = None
        self.model_name = "text-embedding-3-small"
        self.dimensions = 384
        self._initialize_client()
    
    def _initialize_client(self):
        """Initialize OpenAI client."""
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY environment variable is required")
        
        self.client = openai.OpenAI(api_key=api_key)
        logger.info("Initialized OpenAI embedding client")
    
    def generate_embedding(self, text: str) -> List[float]:
        """Generate embedding for text."""
        try:
            # Truncate text if too long
            max_tokens = 8000  # Conservative limit
            if len(text) > max_tokens:
                text = text[:max_tokens] + "..."
            
            response = self.client.embeddings.create(
                model=self.model_name,
                input=text,
                dimensions=self.dimensions
            )
            
            return response.data[0].embedding
        except Exception as e:
            logger.error(f"Error generating embedding: {e}")
            raise

class QdrantManager:
    """Manages Qdrant vector storage for instruction texts."""
    
    def __init__(self):
        self.client = None
        self.collection_name = "instruction_texts"
        self.vector_size = 384
        self._initialize_client()
    
    def _initialize_client(self):
        """Initialize Qdrant client."""
        qdrant_url = os.getenv("QDRANT_URL")
        qdrant_api_key = os.getenv("QDRANT_API_KEY")
        
        if qdrant_url and qdrant_api_key:
            self.client = QdrantClient(url=qdrant_url, api_key=qdrant_api_key)
            logger.info("Initialized Qdrant Cloud client")
        else:
            # Try local Qdrant
            try:
                self.client = QdrantClient("localhost", port=6333)
                logger.info("Initialized local Qdrant client")
            except Exception as e:
                logger.error(f"Failed to connect to Qdrant: {e}")
                raise
        
        self._ensure_collection_exists()
    
    def _ensure_collection_exists(self):
        """Ensure the instruction_texts collection exists."""
        try:
            collections = self.client.get_collections()
            collection_names = [col.name for col in collections.collections]
            
            if self.collection_name not in collection_names:
                self.client.create_collection(
                    collection_name=self.collection_name,
                    vectors_config=VectorParams(
                        size=self.vector_size,
                        distance=Distance.COSINE
                    )
                )
                logger.info(f"Created Qdrant collection: {self.collection_name}")
            else:
                logger.info(f"Qdrant collection exists: {self.collection_name}")
        except Exception as e:
            logger.error(f"Error ensuring collection exists: {e}")
            raise
    
    def upsert_document(self, file_path: Path, text: str, embedding: List[float]) -> bool:
        """Upsert document embedding into Qdrant."""
        try:
            # Generate unique ID based on file path
            doc_id = abs(hash(str(file_path))) % (10**9)
            
            # Prepare metadata
            payload = {
                "filename": file_path.name,
                "filepath": str(file_path),
                "text": text[:1000] + "..." if len(text) > 1000 else text,  # Truncate for storage
                "full_text": text,
                "file_size": file_path.stat().st_size,
                "modified_time": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                "processed_time": datetime.utcnow().isoformat()
            }
            
            point = PointStruct(
                id=doc_id,
                vector=embedding,
                payload=payload
            )
            
            self.client.upsert(
                collection_name=self.collection_name,
                points=[point]
            )
            
            logger.info(f"Upserted document: {file_path.name}")
            return True
        except Exception as e:
            logger.error(f"Error upserting document {file_path.name}: {e}")
            return False
    
    def search_similar_documents(self, query_embedding: List[float], top_k: int = 3) -> List[Dict[str, Any]]:
        """Search for similar documents."""
        try:
            search_result = self.client.search(
                collection_name=self.collection_name,
                query_vector=query_embedding,
                limit=top_k,
                score_threshold=0.5
            )
            
            results = []
            for hit in search_result:
                result = {
                    "filename": hit.payload.get("filename"),
                    "filepath": hit.payload.get("filepath"),
                    "text_excerpt": hit.payload.get("text"),
                    "full_text": hit.payload.get("full_text"),
                    "score": float(hit.score),
                    "modified_time": hit.payload.get("modified_time"),
                    "processed_time": hit.payload.get("processed_time")
                }
                results.append(result)
            
            return results
        except Exception as e:
            logger.error(f"Error searching documents: {e}")
            return []

class InstructionProcessor:
    """Main processor for instruction documents."""
    
    def __init__(self, instructions_dir: str = "instructions"):
        self.instructions_dir = Path(instructions_dir)
        self.converter = DocumentConverter()
        self.embedding_generator = EmbeddingGenerator()
        self.qdrant_manager = QdrantManager()
        
        self.processed_files = []
        self.failed_files = []
    
    def scan_and_process_all(self) -> Dict[str, Any]:
        """Scan instructions directory and process all supported files."""
        if not self.instructions_dir.exists():
            logger.warning(f"Instructions directory not found: {self.instructions_dir}")
            self.instructions_dir.mkdir(parents=True, exist_ok=True)
            self._create_sample_instructions()
        
        # Find all supported files
        supported_files = []
        for extension in self.converter.supported_extensions:
            supported_files.extend(self.instructions_dir.glob(f"**/*{extension}"))
        
        logger.info(f"Found {len(supported_files)} supported files")
        
        # Process each file
        total_processed = 0
        total_failed = 0
        
        for file_path in supported_files:
            try:
                self._process_single_file(file_path)
                self.processed_files.append(file_path)
                total_processed += 1
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {e}")
                self.failed_files.append((file_path, str(e)))
                total_failed += 1
        
        # Return summary
        return {
            "total_files": len(supported_files),
            "processed": total_processed,
            "failed": total_failed,
            "processed_files": [str(f) for f in self.processed_files],
            "failed_files": [(str(f), err) for f, err in self.failed_files]
        }
    
    def _process_single_file(self, file_path: Path):
        """Process a single instruction file."""
        logger.info(f"Processing: {file_path.name}")
        
        # Convert to Markdown
        markdown_text = self.converter.convert_to_markdown(file_path)
        
        # Generate embedding
        embedding = self.embedding_generator.generate_embedding(markdown_text)
        
        # Store in Qdrant
        success = self.qdrant_manager.upsert_document(file_path, markdown_text, embedding)
        
        if not success:
            raise RuntimeError(f"Failed to store in Qdrant: {file_path}")
    
    def search_instructions(self, query: str, top_k: int = 3) -> List[Dict[str, Any]]:
        """Search for relevant instructions based on a query."""
        try:
            # Generate embedding for query
            query_embedding = self.embedding_generator.generate_embedding(query)
            
            # Search in Qdrant
            results = self.qdrant_manager.search_similar_documents(query_embedding, top_k)
            
            logger.info(f"Found {len(results)} relevant instructions for query: {query[:50]}...")
            return results
        except Exception as e:
            logger.error(f"Error searching instructions: {e}")
            return []
    
    def _create_sample_instructions(self):
        """Create sample instruction files for testing."""
        sample_files = [
            {
                "filename": "authentication_guide.txt",
                "content": """# Authentication Guide

## Login Issues
When users experience login problems:

1. Verify username and password
2. Check account status (locked/suspended)
3. Verify two-factor authentication
4. Clear browser cache
5. Try incognito mode
6. Check for service outages

## Common Error Codes
- AUTH001: Invalid credentials
- AUTH002: Account locked
- AUTH003: 2FA verification failed
- AUTH004: Session expired

## Troubleshooting Steps
For persistent login issues:
1. Reset password
2. Disable/re-enable 2FA
3. Contact system administrator
"""
            },
            {
                "filename": "billing_procedures.txt",
                "content": """# Billing Procedures

## Payment Processing
For payment-related inquiries:

1. Verify payment method validity
2. Check sufficient funds/credit
3. Confirm billing address
4. Review transaction history
5. Check for declined transactions

## Subscription Management
- Upgrade/downgrade plans
- Cancel subscriptions
- Refund processing
- Proration calculations

## Error Resolution
Common billing errors:
- PAY001: Card declined
- PAY002: Insufficient funds
- PAY003: Invalid payment method
- PAY004: Billing address mismatch
"""
            },
            {
                "filename": "api_documentation.txt",
                "content": """# API Documentation

## Rate Limiting
API rate limits by tier:
- Basic: 100 requests/hour
- Pro: 1000 requests/hour  
- Enterprise: 10000 requests/hour

## Authentication
Use Bearer tokens in Authorization header:
```
Authorization: Bearer your_api_key_here
```

## Common Issues
- Rate limit exceeded (429 error)
- Invalid API key (401 error)
- Malformed requests (400 error)
- Server errors (500 error)

## Best Practices
1. Implement exponential backoff
2. Cache responses when possible
3. Use webhooks for real-time updates
4. Monitor rate limit headers
"""
            }
        ]
        
        for sample in sample_files:
            file_path = self.instructions_dir / sample["filename"]
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(sample["content"])
        
        logger.info(f"Created {len(sample_files)} sample instruction files")

def main():
    """Main function to run instruction processing."""
    try:
        # Initialize processor
        processor = InstructionProcessor()
        
        # Process all files
        logger.info("Starting instruction file processing...")
        results = processor.scan_and_process_all()
        
        # Print summary
        logger.info("\n" + "="*60)
        logger.info("INSTRUCTION PROCESSING SUMMARY")
        logger.info("="*60)
        logger.info(f"Total files found: {results['total_files']}")
        logger.info(f"Successfully processed: {results['processed']}")
        logger.info(f"Failed to process: {results['failed']}")
        
        if results['processed_files']:
            logger.info("\nProcessed files:")
            for file_path in results['processed_files']:
                logger.info(f"  ✓ {file_path}")
        
        if results['failed_files']:
            logger.info("\nFailed files:")
            for file_path, error in results['failed_files']:
                logger.info(f"  ✗ {file_path}: {error}")
        
        # Test search functionality
        if results['processed'] > 0:
            logger.info("\n" + "="*60)
            logger.info("TESTING SEARCH FUNCTIONALITY")
            logger.info("="*60)
            
            test_queries = [
                "How to fix login problems?",
                "Payment processing issues",
                "API rate limits"
            ]
            
            for query in test_queries:
                logger.info(f"\nQuery: {query}")
                search_results = processor.search_instructions(query, top_k=2)
                
                if search_results:
                    for i, result in enumerate(search_results, 1):
                        logger.info(f"  {i}. {result['filename']} (score: {result['score']:.3f})")
                        logger.info(f"     Excerpt: {result['text_excerpt'][:100]}...")
                else:
                    logger.info("  No relevant results found")
        
        return results['failed'] == 0
        
    except Exception as e:
        logger.error(f"Error in main processing: {e}")
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)