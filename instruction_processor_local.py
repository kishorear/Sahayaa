"""
Instruction Document Processor with local fallback storage
Processes instruction files and provides search functionality without requiring Qdrant.
"""

import os
import sys
import json
import logging
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
import pickle
from datetime import datetime
import numpy as np

# Document processing imports
try:
    from docx import Document
    from openpyxl import load_workbook
    from PyPDF2 import PdfReader
    from pptx import Presentation
    OFFICE_SUPPORT = True
except ImportError as e:
    logging.warning(f"Office document support limited: {e}")
    OFFICE_SUPPORT = False

import openai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class DocumentConverter:
    """Converts various document formats to Markdown text."""
    
    def __init__(self):
        self.supported_extensions = {'.txt', '.pdf', '.docx', '.pptx', '.xlsx'}
    
    def convert_to_markdown(self, file_path: Path) -> str:
        """Convert document to Markdown format."""
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.txt':
                return self._convert_txt(file_path)
            elif extension == '.pdf':
                return self._convert_pdf(file_path)
            elif extension == '.docx':
                return self._convert_docx(file_path)
            elif extension == '.pptx':
                return self._convert_pptx(file_path)
            elif extension == '.xlsx':
                return self._convert_xlsx(file_path)
            else:
                raise ValueError(f"Unsupported file format: {extension}")
        except Exception as e:
            logger.error(f"Error converting {file_path}: {e}")
            raise
    
    def _convert_txt(self, file_path: Path) -> str:
        """Convert plain text file to Markdown."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read().strip()
        
        filename = file_path.stem
        return f"# {filename}\n\n{content}"
    
    def _convert_pdf(self, file_path: Path) -> str:
        """Convert PDF to Markdown."""
        if not OFFICE_SUPPORT:
            raise ImportError("PDF processing requires PyPDF2")
        
        try:
            reader = PdfReader(str(file_path))
            text_content = []
            
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    text_content.append(f"## Page {i + 1}\n\n{page_text}")
            
            filename = file_path.stem
            content = "\n\n".join(text_content)
            return f"# {filename}\n\n{content}"
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            raise
    
    def _convert_docx(self, file_path: Path) -> str:
        """Convert Word document to Markdown."""
        if not OFFICE_SUPPORT:
            raise ImportError("DOCX processing requires python-docx")
        
        try:
            doc = Document(str(file_path))
            content_parts = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    if paragraph.style.name.startswith('Heading'):
                        level = min(int(paragraph.style.name.split()[-1]) if paragraph.style.name.split()[-1].isdigit() else 1, 6)
                        content_parts.append(f"{'#' * level} {text}")
                    else:
                        content_parts.append(text)
            
            for table in doc.tables:
                table_md = self._table_to_markdown(table)
                if table_md:
                    content_parts.append(table_md)
            
            filename = file_path.stem
            content = "\n\n".join(content_parts)
            return f"# {filename}\n\n{content}"
        except Exception as e:
            logger.error(f"Error reading DOCX {file_path}: {e}")
            raise
    
    def _convert_pptx(self, file_path: Path) -> str:
        """Convert PowerPoint presentation to Markdown."""
        if not OFFICE_SUPPORT:
            raise ImportError("PPTX processing requires python-pptx")
        
        try:
            prs = Presentation(str(file_path))
            slides_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slide_content = "\n".join(slide_text)
                    slides_content.append(f"## Slide {i + 1}\n\n{slide_content}")
            
            filename = file_path.stem
            content = "\n\n".join(slides_content)
            return f"# {filename}\n\n{content}"
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path}: {e}")
            raise
    
    def _convert_xlsx(self, file_path: Path) -> str:
        """Convert Excel spreadsheet to Markdown."""
        if not OFFICE_SUPPORT:
            raise ImportError("XLSX processing requires openpyxl")
        
        try:
            workbook = load_workbook(str(file_path), data_only=True)
            sheets_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    if any(cell for cell in row if cell is not None):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        sheet_data.append("| " + " | ".join(row_data) + " |")
                
                if sheet_data:
                    if len(sheet_data) > 1:
                        header_separator = "| " + " | ".join(["---"] * len(sheet_data[0].split("|")[1:-1])) + " |"
                        sheet_data.insert(1, header_separator)
                    
                    sheet_content = f"## {sheet_name}\n\n" + "\n".join(sheet_data)
                    sheets_content.append(sheet_content)
            
            filename = file_path.stem
            content = "\n\n".join(sheets_content)
            return f"# {filename}\n\n{content}"
        except Exception as e:
            logger.error(f"Error reading XLSX {file_path}: {e}")
            raise
    
    def _table_to_markdown(self, table) -> str:
        """Convert Word table to Markdown format."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append("| " + " | ".join(cells) + " |")
        
        if len(rows) > 1:
            header_separator = "| " + " | ".join(["---"] * len(rows[0].split("|")[1:-1])) + " |"
            rows.insert(1, header_separator)
        
        return "\n".join(rows) if rows else ""

class EmbeddingGenerator:
    """Generates embeddings using OpenAI."""
    
    def __init__(self):
        self.client = None
        self.model_name = "text-embedding-3-small"
        self.dimensions = 384
        self._initialize_client()
    
    def _initialize_client(self):
        """Initialize OpenAI client."""
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY environment variable is required")
        
        self.client = openai.OpenAI(api_key=api_key)
        logger.info("Initialized OpenAI embedding client")
    
    def generate_embedding(self, text: str) -> List[float]:
        """Generate embedding for text."""
        try:
            max_tokens = 8000
            if len(text) > max_tokens:
                text = text[:max_tokens] + "..."
            
            response = self.client.embeddings.create(
                model=self.model_name,
                input=text,
                dimensions=self.dimensions
            )
            
            return response.data[0].embedding
        except Exception as e:
            logger.error(f"Error generating embedding: {e}")
            raise

class LocalVectorStorage:
    """Local file-based vector storage as fallback."""
    
    def __init__(self, storage_dir: str = "vector_storage"):
        self.storage_dir = Path(storage_dir)
        self.storage_dir.mkdir(exist_ok=True)
        self.index_file = self.storage_dir / "instruction_index.json"
        self.embeddings_file = self.storage_dir / "embeddings.pkl"
        self.documents = {}
        self.embeddings = {}
        self._load_storage()
    
    def _load_storage(self):
        """Load existing storage."""
        if self.index_file.exists():
            with open(self.index_file, 'r') as f:
                self.documents = json.load(f)
        
        if self.embeddings_file.exists():
            with open(self.embeddings_file, 'rb') as f:
                self.embeddings = pickle.load(f)
        
        logger.info(f"Loaded {len(self.documents)} documents from local storage")
    
    def _save_storage(self):
        """Save storage to files."""
        with open(self.index_file, 'w') as f:
            json.dump(self.documents, f, indent=2, default=str)
        
        with open(self.embeddings_file, 'wb') as f:
            pickle.dump(self.embeddings, f)
    
    def store_document(self, file_path: Path, text: str, embedding: List[float]) -> bool:
        """Store document and embedding."""
        try:
            doc_id = str(file_path)
            
            self.documents[doc_id] = {
                "filename": file_path.name,
                "filepath": str(file_path),
                "text": text[:1000] + "..." if len(text) > 1000 else text,
                "full_text": text,
                "file_size": file_path.stat().st_size,
                "modified_time": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                "processed_time": datetime.utcnow().isoformat()
            }
            
            self.embeddings[doc_id] = embedding
            self._save_storage()
            
            logger.info(f"Stored document: {file_path.name}")
            return True
        except Exception as e:
            logger.error(f"Error storing document {file_path.name}: {e}")
            return False
    
    def search_similar_documents(self, query_embedding: List[float], top_k: int = 3) -> List[Dict[str, Any]]:
        """Search for similar documents using cosine similarity."""
        try:
            if not self.embeddings:
                return []
            
            query_vector = np.array(query_embedding)
            similarities = []
            
            for doc_id, doc_embedding in self.embeddings.items():
                doc_vector = np.array(doc_embedding)
                
                # Calculate cosine similarity
                dot_product = np.dot(query_vector, doc_vector)
                norm_query = np.linalg.norm(query_vector)
                norm_doc = np.linalg.norm(doc_vector)
                
                if norm_query > 0 and norm_doc > 0:
                    similarity = dot_product / (norm_query * norm_doc)
                    similarities.append((doc_id, similarity))
            
            # Sort by similarity and take top_k
            similarities.sort(key=lambda x: x[1], reverse=True)
            top_results = similarities[:top_k]
            
            # Filter by minimum threshold
            results = []
            for doc_id, score in top_results:
                if score > 0.5:  # Minimum similarity threshold
                    doc_info = self.documents.get(doc_id, {})
                    result = {
                        "filename": doc_info.get("filename"),
                        "filepath": doc_info.get("filepath"),
                        "text_excerpt": doc_info.get("text"),
                        "full_text": doc_info.get("full_text"),
                        "score": float(score),
                        "modified_time": doc_info.get("modified_time"),
                        "processed_time": doc_info.get("processed_time")
                    }
                    results.append(result)
            
            return results
        except Exception as e:
            logger.error(f"Error searching documents: {e}")
            return []

class InstructionProcessor:
    """Main processor for instruction documents."""
    
    def __init__(self, instructions_dir: str = "instructions"):
        self.instructions_dir = Path(instructions_dir)
        self.converter = DocumentConverter()
        self.embedding_generator = EmbeddingGenerator()
        self.vector_storage = LocalVectorStorage()
        
        self.processed_files = []
        self.failed_files = []
    
    def scan_and_process_all(self) -> Dict[str, Any]:
        """Scan instructions directory and process all supported files."""
        if not self.instructions_dir.exists():
            logger.warning(f"Instructions directory not found: {self.instructions_dir}")
            self.instructions_dir.mkdir(parents=True, exist_ok=True)
            self._create_sample_instructions()
        
        # Find all supported files
        supported_files = []
        for extension in self.converter.supported_extensions:
            supported_files.extend(self.instructions_dir.glob(f"**/*{extension}"))
        
        logger.info(f"Found {len(supported_files)} supported files")
        
        # Process each file
        total_processed = 0
        total_failed = 0
        
        for file_path in supported_files:
            try:
                self._process_single_file(file_path)
                self.processed_files.append(file_path)
                total_processed += 1
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {e}")
                self.failed_files.append((file_path, str(e)))
                total_failed += 1
        
        return {
            "total_files": len(supported_files),
            "processed": total_processed,
            "failed": total_failed,
            "processed_files": [str(f) for f in self.processed_files],
            "failed_files": [(str(f), err) for f, err in self.failed_files]
        }
    
    def _process_single_file(self, file_path: Path):
        """Process a single instruction file."""
        logger.info(f"Processing: {file_path.name}")
        
        # Convert to Markdown
        markdown_text = self.converter.convert_to_markdown(file_path)
        
        # Generate embedding
        embedding = self.embedding_generator.generate_embedding(markdown_text)
        
        # Store in local vector storage
        success = self.vector_storage.store_document(file_path, markdown_text, embedding)
        
        if not success:
            raise RuntimeError(f"Failed to store document: {file_path}")
    
    def search_instructions(self, query: str, top_k: int = 3) -> List[Dict[str, Any]]:
        """Search for relevant instructions based on a query."""
        try:
            # Generate embedding for query
            query_embedding = self.embedding_generator.generate_embedding(query)
            
            # Search in local storage
            results = self.vector_storage.search_similar_documents(query_embedding, top_k)
            
            logger.info(f"Found {len(results)} relevant instructions for query: {query[:50]}...")
            return results
        except Exception as e:
            logger.error(f"Error searching instructions: {e}")
            return []
    
    def _create_sample_instructions(self):
        """Create sample instruction files for testing."""
        sample_files = [
            {
                "filename": "authentication_guide.txt",
                "content": """# Authentication Guide

## Login Issues
When users experience login problems:

1. Verify username and password
2. Check account status (locked/suspended)
3. Verify two-factor authentication
4. Clear browser cache
5. Try incognito mode
6. Check for service outages

## Common Error Codes
- AUTH001: Invalid credentials
- AUTH002: Account locked
- AUTH003: 2FA verification failed
- AUTH004: Session expired

## Troubleshooting Steps
For persistent login issues:
1. Reset password
2. Disable/re-enable 2FA
3. Contact system administrator
"""
            },
            {
                "filename": "billing_procedures.txt",
                "content": """# Billing Procedures

## Payment Processing
For payment-related inquiries:

1. Verify payment method validity
2. Check sufficient funds/credit
3. Confirm billing address
4. Review transaction history
5. Check for declined transactions

## Subscription Management
- Upgrade/downgrade plans
- Cancel subscriptions
- Refund processing
- Proration calculations

## Error Resolution
Common billing errors:
- PAY001: Card declined
- PAY002: Insufficient funds
- PAY003: Invalid payment method
- PAY004: Billing address mismatch
"""
            },
            {
                "filename": "api_documentation.txt",
                "content": """# API Documentation

## Rate Limiting
API rate limits by tier:
- Basic: 100 requests/hour
- Pro: 1000 requests/hour  
- Enterprise: 10000 requests/hour

## Authentication
Use Bearer tokens in Authorization header:
```
Authorization: Bearer your_api_key_here
```

## Common Issues
- Rate limit exceeded (429 error)
- Invalid API key (401 error)
- Malformed requests (400 error)
- Server errors (500 error)

## Best Practices
1. Implement exponential backoff
2. Cache responses when possible
3. Use webhooks for real-time updates
4. Monitor rate limit headers
"""
            }
        ]
        
        for sample in sample_files:
            file_path = self.instructions_dir / sample["filename"]
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(sample["content"])
        
        logger.info(f"Created {len(sample_files)} sample instruction files")

def main():
    """Main function to run instruction processing."""
    try:
        processor = InstructionProcessor()
        
        logger.info("Starting instruction file processing...")
        results = processor.scan_and_process_all()
        
        logger.info("\n" + "="*60)
        logger.info("INSTRUCTION PROCESSING SUMMARY")
        logger.info("="*60)
        logger.info(f"Total files found: {results['total_files']}")
        logger.info(f"Successfully processed: {results['processed']}")
        logger.info(f"Failed to process: {results['failed']}")
        
        if results['processed_files']:
            logger.info("\nProcessed files:")
            for file_path in results['processed_files']:
                logger.info(f"  ✓ {file_path}")
        
        if results['failed_files']:
            logger.info("\nFailed files:")
            for file_path, error in results['failed_files']:
                logger.info(f"  ✗ {file_path}: {error}")
        
        # Test search functionality
        if results['processed'] > 0:
            logger.info("\n" + "="*60)
            logger.info("TESTING SEARCH FUNCTIONALITY")
            logger.info("="*60)
            
            test_queries = [
                "How to fix login problems?",
                "Payment processing issues",
                "API rate limits"
            ]
            
            for query in test_queries:
                logger.info(f"\nQuery: {query}")
                search_results = processor.search_instructions(query, top_k=2)
                
                if search_results:
                    for i, result in enumerate(search_results, 1):
                        logger.info(f"  {i}. {result['filename']} (score: {result['score']:.3f})")
                        logger.info(f"     Excerpt: {result['text_excerpt'][:100]}...")
                else:
                    logger.info("  No relevant results found")
        
        return results['failed'] == 0
        
    except Exception as e:
        logger.error(f"Error in main processing: {e}")
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)